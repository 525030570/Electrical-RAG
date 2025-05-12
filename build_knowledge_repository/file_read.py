import os
from datetime import datetime
from pdfminer.high_level import extract_text as pdf_extract_text
import docx2txt
from pptx import Presentation
from llama_index.core.schema import Document

# --------------------------
# 文件读取函数定义（封装为 LlamaIndex Document）
# --------------------------

def read_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
    except Exception as e:
        return Document(
            text=f"读取 TXT 文件失败: {str(e)}",
            metadata={"error": True}
        )
    return Document(
        text=content,
        metadata={
            "filename": os.path.basename(file_path),
            "file_path": file_path,
            "file_type": ".txt",
            "read_time": datetime.now().isoformat(),
        }
    )

def read_docx(file_path):
    try:
        content = docx2txt.process(file_path)
    except Exception as e:
        return Document(
            text=f"读取 DOCX 文件失败: {str(e)}",
            metadata={"error": True}
        )
    return Document(
        text=content,
        metadata={
            "filename": os.path.basename(file_path),
            "file_path": file_path,
            "file_type": ".docx",
            "read_time": datetime.now().isoformat(),
        }
    )

def read_pdf(file_path):
    try:
        content = pdf_extract_text(file_path)
    except Exception as e:
        return Document(
            text=f"读取 PDF 文件失败: {str(e)}",
            metadata={"error": True}
        )
    return Document(
        text=content,
        metadata={
            "filename": os.path.basename(file_path),
            "file_path": file_path,
            "file_type": ".pdf",
            "read_time": datetime.now().isoformat(),
        }
    )

def read_pptx(file_path):
    try:
        presentation = Presentation(file_path)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        content = "\n".join(full_text)
    except Exception as e:
        return Document(
            text=f"读取 PPTX 文件失败: {str(e)}",
            metadata={"error": True}
        )
    return Document(
        text=content,
        metadata={
            "filename": os.path.basename(file_path),
            "file_path": file_path,
            "file_type": ".pptx",
            "read_time": datetime.now().isoformat(),
        }
    )

# --------------------------
# 通用文件读取函数（返回 LlamaIndex Document）
# --------------------------

def read_file(file_path):
    """
    通用读取函数，支持：.txt, .docx, .pdf, .pptx
    返回一个 LlamaIndex Document 对象
    """
    if not os.path.isfile(file_path):
        return Document(
            text=f"文件不存在: {file_path}",
            metadata={"error": True}
        )

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".txt":
            return read_txt(file_path)
        elif ext == ".docx":
            return read_docx(file_path)
        elif ext == ".pdf":
            return read_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            return read_pptx(file_path)
        else:
            return Document(
                text=f"不支持的文件格式: {ext}",
                metadata={"error": True}
            )
    except Exception as e:
        return Document(
            text=f"读取文件失败: {str(e)}",
            metadata={"error": True}
        )